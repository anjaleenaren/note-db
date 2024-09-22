import streamlit as st
from langchain_openai.chat_models import ChatOpenAI
import os
from dotenv import load_dotenv
import sqlite3
import datetime
import boto3
from botocore.exceptions import NoCredentialsError
import uuid
import pytesseract
from PIL import Image
import io
from PyPDF2 import PdfReader
from pdf2image import convert_from_bytes
from pptx import Presentation
import re
import whisper
import tempfile
from functools import lru_cache
from ai_ta import AI_TA
import openai
import numpy as np
from pydub import AudioSegment
import time

load_dotenv()

# Database setup
@st.cache_resource
def init_db():
    conn = sqlite3.connect('class_notes.db', check_same_thread=False)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS classes
                 (id INTEGER PRIMARY KEY, 
                  name TEXT UNIQUE)''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS ta_indexes
                 (id INTEGER PRIMARY KEY, 
                  class_name TEXT UNIQUE,
                  serialized_index TEXT UNIQUE,
                  FOREIGN KEY (class_name) REFERENCES classes(name))''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS notes
                 (id INTEGER PRIMARY KEY, class_id INTEGER, title TEXT, content TEXT, timestamp DATETIME, file_urls TEXT, audio_urls TEXT,
                 FOREIGN KEY (class_id) REFERENCES classes(id))''')
    
    # Check if file_urls and audio_urls columns exist, if not, add them
    c.execute("PRAGMA table_info(notes)")
    columns = [column[1] for column in c.fetchall()]
    if 'file_urls' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN file_urls TEXT")
    if 'audio_urls' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN audio_urls TEXT")
    
    conn.commit()
    return conn

if 'db_conn' not in st.session_state:
    st.session_state.db_conn = init_db()

# Helper functions
def add_class(class_name):
    c = st.session_state.db_conn.cursor()
    c.execute("INSERT OR IGNORE INTO classes (name) VALUES (?)", (class_name,))

    class_ta = AI_TA(class_name)
    c.execute("INSERT OR IGNORE INTO ta_indexes (class_name, serialized_index) VALUES (?, ?)", 
                (class_name, class_ta.serialize()))
    st.session_state.db_conn.commit()

def get_classes():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT name FROM classes")
    return [row[0] for row in c.fetchall()]

def extract_text_from_file(file):
    text = ""
    if file.type == "application/pdf":
        pdf_reader = PdfReader(file)
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text.strip():
                text += page_text + "\n"
            else:
                # If no text was extracted, the page might be scanned. Use OCR.
                file.seek(0)  # Reset file pointer to the beginning
                images = convert_from_bytes(file.read(), first_page=page_num+1, last_page=page_num+1)
                for image in images:
                    text += pytesseract.image_to_string(image) + "\n"
    elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
        prs = Presentation(io.BytesIO(file.read()))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
            text += "\n"
    else:
        # Process as image
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
    
    # Strip multiple consecutive new lines to be one new line at most
    text = re.sub(r'\n{2,}', '\n', text)
    return text.strip()

def add_note(class_name, note_content, note_title, files=None, audio_files=None):
    print("Add note id " + note_title)
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT id FROM classes WHERE name = ?", (class_name,))
    class_id = c.fetchone()[0]
    
    file_urls = []
    audio_urls = []
    
    if files:
        for file in files:
            file_extension = os.path.splitext(file.name)[1]
            file_filename = f"files/{uuid.uuid4()}{file_extension}"
            if upload_to_s3(file, os.getenv('S3_BUCKET_NAME'), file_filename):
                file_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), file_filename))
    
    if audio_files:
        for audio_file in audio_files:
            audio_filename = f"audio/{uuid.uuid4()}{os.path.splitext(audio_file.name)[1]}"
            if upload_to_s3(audio_file, os.getenv('S3_BUCKET_NAME'), audio_filename):
                audio_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), audio_filename))
    
    # Strip multiple consecutive new lines in note_content
    note_content = re.sub(r'\n{2,}', '\n', note_content)
    
    c.execute("""INSERT INTO notes (class_id, title, content, timestamp, file_urls, audio_urls) 
                 VALUES (?, ?, ?, ?, ?, ?)""",
              (class_id, note_title, note_content, datetime.datetime.now(), ','.join(file_urls), ','.join(audio_urls)))
    
    new_note_id = c.lastrowid  # Get the id of the newly inserted note
 
    return new_note_id  # Return the id of the newly inserted note

def train_ai_ta(class_name, note_content, title = "title"):
    c = st.session_state.db_conn.cursor()

    # Get the serialized index of the AI-TA for the class
    c.execute("SELECT serialized_index FROM ta_indexes WHERE class_name = ?", (class_name,))
    serialized_index = c.fetchone()[0]
    class_ta = AI_TA.deserialize(serialized_index)

    #Train the AI-TA with the new note content
    class_ta.train(note_content, title)

    # Update the serialized index in the database
    c.execute("UPDATE ta_indexes SET serialized_index = ? WHERE class_name = ?", (class_ta.serialize(), class_name))
    st.session_state.db_conn.commit()

def update_note(note_id, note_content, note_title, files=None, audio_files=None):
    print(f"Update note id {note_id} {note_title}")
    c = st.session_state.db_conn.cursor()
    
    file_urls = []
    audio_urls = []
    
    if files:
        for file in files:
            file_extension = os.path.splitext(file.name)[1]
            file_filename = f"files/{uuid.uuid4()}{file_extension}"
            if upload_to_s3(file, os.getenv('S3_BUCKET_NAME'), file_filename):
                file_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), file_filename))
    
    if audio_files:
        for audio_file in audio_files:
            audio_filename = f"audio/{uuid.uuid4()}{os.path.splitext(audio_file.name)[1]}"
            if upload_to_s3(audio_file, os.getenv('S3_BUCKET_NAME'), audio_filename):
                audio_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), audio_filename))
    
    # Get existing file_urls and audio_urls
    c.execute("SELECT file_urls, audio_urls FROM notes WHERE id = ?", (note_id,))
    existing_file_urls, existing_audio_urls = c.fetchone()
    if existing_file_urls:
        file_urls = existing_file_urls.split(',') + file_urls
    if existing_audio_urls:
        audio_urls = existing_audio_urls.split(',') + audio_urls
    
    # Strip multiple consecutive new lines in content
    note_content = re.sub(r'\n{2,}', '\n', note_content)
    
    c.execute("""UPDATE notes 
                 SET title = ?, content = ?, timestamp = ?, file_urls = ?, audio_urls = ?
                 WHERE id = ?""",
              (note_title, note_content, datetime.datetime.now(), ','.join(file_urls), ','.join(audio_urls), note_id))
    st.session_state.db_conn.commit()

def delete_file_from_note(note_id, file_url):
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT file_urls FROM notes WHERE id = ?", (note_id,))
    file_urls = c.fetchone()[0].split(',')
    file_urls.remove(file_url)
    c.execute("UPDATE notes SET file_urls = ? WHERE id = ?", (','.join(file_urls), note_id))

def get_notes(class_name):
    c = st.session_state.db_conn.cursor()
    try:
        c.execute("""SELECT notes.id, notes.content, notes.timestamp, notes.file_urls, notes.audio_urls, notes.title
                     FROM notes 
                     JOIN classes ON notes.class_id = classes.id 
                     WHERE classes.name = ?
                     ORDER BY notes.timestamp DESC""", (class_name,))
    except sqlite3.OperationalError:
        # If the query fails, fall back to the original schema
        c.execute("""SELECT notes.id, notes.content, notes.timestamp, NULL as file_urls, NULL as audio_urls, notes.title
                     FROM notes 
                     JOIN classes ON notes.class_id = classes.id 
                     WHERE classes.name = ?
                     ORDER BY notes.timestamp DESC""", (class_name,))
    return c.fetchall()

def get_tables():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = c.fetchall()
    return [table[0] for table in tables]

# Add these new functions
def get_classes_data():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT * FROM classes")
    return c.fetchall()

def get_notes_data():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT * FROM notes")
    return c.fetchall()

# Add this new function
def clear_database():
    c = st.session_state.db_conn.cursor()
    c.execute("DELETE FROM notes")
    c.execute("DELETE FROM classes")
    # c.execute("DELETE FROM sqlite_sequence WHERE name IN ('notes', 'classes')")  # Reset auto-increment
    st.session_state.db_conn.commit()


# Add these new functions
def save_note(selected_class, note_id, note_content, note_title, files=None, audio_files=None):
    if not note_title:
        note_title = generate_title(note_content)
    print("Note title = " + note_title)
    if note_id:
        update_note(note_id, note_content, note_title, files, audio_files)
        st.success("Note updated!")
    else:
        new_note_id = add_note(selected_class, note_content, note_title, files, audio_files)
        st.session_state.current_note_id = new_note_id  # Save the new note id
        st.success("New note added!")
    return (note_id or new_note_id, note_title)

def get_note_by_id(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT title, content FROM notes WHERE id = ?", (note_id,))
    return c.fetchone()

# Add this new function to delete a note
def delete_note(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("DELETE FROM notes WHERE id = ?", (note_id,))
    st.session_state.db_conn.commit()

# Add these new functions for S3 operations
def upload_to_s3(file, bucket_name, object_name):
    pass
    s3_client = boto3.client('s3',
                             aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
                             aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'))
    try:
        s3_client.upload_fileobj(file, bucket_name, object_name)
        return True
    except NoCredentialsError:
        st.error("AWS credentials not available")
        return False
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return False

def get_s3_url(bucket_name, object_name):
    return f"https://{bucket_name}.s3.amazonaws.com/{object_name}"

# Add this cached function for audio transcription
@lru_cache(maxsize=None)
def cached_transcribe_audio(file_content):
    model = whisper.load_model("base")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
        tmp_file.write(file_content)
        tmp_file_path = tmp_file.name

    result = model.transcribe(tmp_file_path)
    os.unlink(tmp_file_path)
    return result["text"]

def generate_title(content):
    if content:
        try:
            model = ChatOpenAI(temperature=0.7, api_key=openai_api_key)
            ret = model.invoke("Generate title for these notes" + content)
            
            return ret.content
        except Exception as e:
            st.error(f"Error generating title: {str(e)}")
            return "Untitled Note"

def stream_transcribe_audio(audio_file, selected_class, note_id, note_title, note_content):
    model = whisper.load_model("base")
    
    # Load the audio file
    audio = AudioSegment.from_file(audio_file)
    
    # Set chunk size (e.g., 10 seconds)
    chunk_duration_ms = 60 * 1000
    
    full_transcript = ""
    placeholder = st.empty()
    
    for i in range(0, len(audio), chunk_duration_ms):
        # Extract a chunk of audio
        chunk = audio[i:i+chunk_duration_ms]
        
        # Export the chunk to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
            chunk.export(tmp_file.name, format="wav")
            tmp_file_path = tmp_file.name
        
        # Transcribe the chunk
        result = model.transcribe(tmp_file_path)
        chunk_transcript = result["text"]
        
        # Append to full transcript and update display
        full_transcript += chunk_transcript + " "
        placeholder.text_area("Transcription in progress:", full_transcript, height=150)
        
        # Append the chunk transcript to the note content and save
        if i == 0:
            updated_content = note_content + f"Audio Transcription: {chunk_transcript}"
        else: 
            updated_content = note_content + f"{chunk_transcript}"
        note_id, note_title = save_note(selected_class, note_id, updated_content, note_title)
        note_content = updated_content  # Update note_content for the next iteration
        
        # Clean up the temporary file
        os.unlink(tmp_file_path)
        
        # Simulate processing time
        time.sleep(0.1)

        # st.rerun()
    
    return full_transcript.strip(), note_id, note_title

# Streamlit UI
st.title("🎓 Class Notes Manager")

# Sidebar for class management
st.sidebar.title("Class Management")

# Display existing classes
st.sidebar.subheader("Existing Classes")
classes = get_classes()
selected_class = st.sidebar.selectbox("Select a class:", [""] + classes)

# Add new class
st.sidebar.subheader("Add a New Class")
new_class = st.sidebar.text_input("Add a new class:")
if st.sidebar.button("Add Class"):
    if new_class:
        add_class(new_class)
        st.sidebar.success(f"Added {new_class}")
        st.rerun()

# OpenAI API Key input
st.sidebar.subheader("Settings")
openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password", value=os.getenv("OPENAI_API_KEY", ""))

# Main content area
if selected_class:
    # AI-TA chat
    st.subheader("AI TA Chat :)")
    if openai_api_key.startswith("sk-"):
        chat_input = st.text_input("Ask AI-TA a question about this class:")
        if st.button("Ask"):
            c = st.session_state.db_conn.cursor()
            c.execute("SELECT serialized_index FROM ta_indexes WHERE class_name = ?", (selected_class,))
            serialized_index = c.fetchone()[0]
            print("serialized_index fetched! ")
            class_ta = AI_TA.deserialize(serialized_index)
            print("class_ta deseralized ", class_ta.class_name)
            openai.api_key = openai_api_key
            response = class_ta.query(chat_input)
            print("response generated ", response)
            st.info(response)
    else:
        st.warning("Please enter your OpenAI API key to use the AI-TA feature!", icon="⚠")

    st.subheader(f"Notes for {selected_class}")
    
    # Initialize session state for current note
    if 'current_note_id' not in st.session_state:
        st.session_state.current_note_id = None
    if 'current_note_content' not in st.session_state:
        st.session_state.current_note_content = ""
    if 'current_note_title' not in st.session_state:
        st.session_state.current_note_title = ""
    if 'transcription_done' not in st.session_state:
        st.session_state.transcription_done = False

    # Text input for notes
    note_title = st.text_input("Note Title:", value=st.session_state.current_note_title)
    note_content = st.text_area("Edit note:", value=st.session_state.current_note_content)
    
    # File upload for multiple images/PDFs/PPTs
    files = st.file_uploader("Upload images, PDFs, or PowerPoint files (text extraction will be performed)", 
                             type=["png", "jpg", "jpeg", "pdf", "ppt", "pptx"], 
                             accept_multiple_files=True)
    if files:
        for i, file in enumerate(files):
            if file.type == "application/pdf":
                st.write(f"PDF uploaded: {file.name}")
            elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                st.write(f"PowerPoint uploaded: {file.name}")
            else:
                st.image(file, caption=f"Uploaded Image: {file.name}", width=200)
            extracted_text = extract_text_from_file(file)
            st.text_area(f"Extracted Text from {file.name}", value=extracted_text, height=150)
            if st.button(f"Append Extracted Text to Note (File {i+1})", key=f"append_button_{i}"):
                note_content += f"\n\nExtracted Text from {file.name}:\n{extracted_text}"
                st.session_state.current_note_content = note_content
                save_note(selected_class, st.session_state.current_note_id, note_content, note_title)
                st.rerun()
    
    # Audio file upload
    audio_files = st.file_uploader("Upload audio files (transcription will be performed)", 
                                   type=["mp3", "wav", "m4a"], 
                                   accept_multiple_files=True)
    
    if 'transcribed_texts' not in st.session_state:
        st.session_state.transcribed_texts = {}

    if audio_files:
        for i, audio_file in enumerate(audio_files):
            st.audio(audio_file)
            
            if not st.session_state.transcription_done:
                with st.spinner(f"Transcribing {audio_file.name}..."):
                    transcribed_text, note_id, note_title = stream_transcribe_audio(audio_file, selected_class, st.session_state.current_note_id, note_title, note_content)
                    st.session_state.current_note_id = note_id
                    st.session_state.current_note_title = note_title
                    st.session_state.current_note_content = note_content  # Update the current note content
                    st.session_state.transcribed_texts[audio_file.name] = transcribed_text
                    st.session_state.transcription_done = True
            
            if audio_file.name in st.session_state.transcribed_texts:
                transcribed_text = st.session_state.transcribed_texts[audio_file.name]
                st.text_area(f"Transcribed Text from {audio_file.name}", value=transcribed_text, height=150)
                if st.button(f"Append Transcribed Text to Note (Audio {i+1})", key=f"append_audio_button_{i}"):
                    note_content += f"\n\nTranscribed Text from {audio_file.name}:\n{transcribed_text}"
                    st.session_state.current_note_content = note_content
                    save_note(selected_class, st.session_state.current_note_id, note_content, note_title)
                    st.rerun()
    
    # Save and New Note buttons
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Save Note"):
            if not st.session_state.current_note_content:
                st.warning("Note content is empty. Please add some content before saving.")
            else: 
                save_note(selected_class, st.session_state.current_note_id, note_content, note_title, files, audio_files)
                st.session_state.current_note_id = None
                st.session_state.current_note_content = ""
                st.session_state.current_note_title = ""
                st.session_state.transcribed_texts = {}
                st.session_state.transcription_done = False  # Reset the transcription flag
            st.rerun()
    
    with col2:
        if st.button("New Note"):
            st.session_state.current_note_id = None
            st.session_state.current_note_content = ""
            st.session_state.current_note_title = ""
            st.session_state.transcribed_texts = {}
            st.session_state.transcription_done = False  # Reset the transcription flag
            # Clear file uploads
            st.session_state.files = None
            st.session_state.audio_files = None
            # Clear extracted text
            if 'extracted_texts' in st.session_state:
                del st.session_state.extracted_texts
            st.rerun()

    # Display existing notes
    st.subheader("Existing Notes")
    for note_id, note_content, timestamp, file_urls, audio_urls, note_title in get_notes(selected_class):
        col1, col2, col3, col4 = st.columns([3, 2, 2, 1])
        with col1:
            date_only = timestamp.split()[0]
            if st.button(f"{date_only}: {note_title} - {note_content[:50]}...", key=f"note_{note_id}"):
                if st.session_state.current_note_id:
                    print("Update note before switching")
                    update_note(st.session_state.current_note_id, note_content, note_title)
                st.session_state.current_note_id = note_id
                title, content = get_note_by_id(note_id)
                st.session_state.current_note_title = title
                st.session_state.current_note_content = content
                st.rerun()
        with col2:
            if file_urls:
                for file_url in file_urls.split(','):
                    if file_url.endswith(".pdf"):
                        st.write("PDF")
                    elif file_url.endswith((".ppt", ".pptx")):
                        st.write("PowerPoint")
                    else:
                        st.image(file_url, width=100)
                    if st.button("Delete File", key=f"delete_file_{note_id}_{file_url}"):
                        delete_file_from_note(note_id, file_url)
                        st.rerun()
        with col3:
            if audio_urls:
                for audio_url in audio_urls.split(','):
                    st.audio(audio_url)
                    if st.button("Delete Audio", key=f"delete_audio_{note_id}_{audio_url}"):
                        delete_file_from_note(note_id, audio_url)
                        st.rerun()
        with col4:
            if st.button("Delete Note", key=f"delete_note_{note_id}"):
                delete_note(note_id)
                if st.session_state.current_note_id == note_id:
                    st.session_state.current_note_id = None
                    st.session_state.current_note_content = ""
                    st.session_state.current_note_title = ""
                st.rerun()
else:
    st.info("Please select a class from the sidebar or add a new one.")

# Close the database connection when the app is done
# st.on_script_run.add(st.session_state.db_conn.close)
